# -*- coding: utf-8 -*-
"""生成 BigAlpha 2026 三大赛事介绍 PPT。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- 主题配色（深蓝 + 强调金） -------------------------------------------------
NAVY = RGBColor(0x0F, 0x2A, 0x5B)        # 主深蓝
BLUE = RGBColor(0x1F, 0x4E, 0xA8)        # 标题蓝
ACCENT = RGBColor(0xE8, 0xB0, 0x3E)      # 金色强调
LIGHT = RGBColor(0xF4, 0xF7, 0xFC)       # 背景浅
GREY = RGBColor(0x4A, 0x55, 0x68)        # 正文灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BLUE = RGBColor(0xDE, 0xE7, 0xF7)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
ORANGE = RGBColor(0xE0, 0x6C, 0x2A)
PURPLE = RGBColor(0x6A, 0x4C, 0x93)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.shadow.inherit = False
    return slide


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=GREY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, *, fill=None, line=None,
             line_width=None, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if line_width is not None:
            sh.line.width = line_width
    return sh


def page_header(slide, title, subtitle=None, idx=None):
    add_rect(slide, 0, 0, SW, Inches(0.55), fill=NAVY)
    add_text(slide, Inches(0.5), Inches(0.08), Inches(11), Inches(0.4),
             title, size=20, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if idx is not None:
        add_text(slide, Inches(11.7), Inches(0.08), Inches(1.4), Inches(0.4),
                 idx, size=11, color=ACCENT, align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.65), Inches(12.5), Inches(0.4),
                 subtitle, size=12, color=BLUE)


def add_bullets(slide, left, top, width, height, items, *,
                size=13, color=GREY, gap=6, bullet="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        if isinstance(item, tuple):
            head, body = item
            r = p.add_run()
            r.text = f"{bullet} {head}"
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = f"  {body}"
            r2.font.name = "Microsoft YaHei"
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = f"{bullet} {item}"
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


# =============================================================================
# Slide 1 — 封面
# =============================================================================
s = add_slide()
add_rect(s, 0, 0, SW, SH, fill=NAVY)
# 装饰条
add_rect(s, 0, Inches(4.6), SW, Inches(0.05), fill=ACCENT)
# 大标题
add_text(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2),
         "BigAlpha 2026 全球高校联赛",
         size=48, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.7),
         "三大赛事赛题整合介绍",
         size=24, color=ACCENT)
add_text(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.5),
         "AI 因子挖掘 · 端到端大模型 · AI 开放创新",
         size=18, color=WHITE)
add_text(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.4),
         "依托 BigQuant AI 投研平台 · 真实数据 · 开放赛题 · 专业评审",
         size=13, color=SOFT_BLUE)
add_text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.4),
         "BigQuant（宽邦科技）  |  2026.05 – 2026.09",
         size=12, color=SOFT_BLUE)

# =============================================================================
# Slide 2 — 目录
# =============================================================================
s = add_slide()
page_header(s, "目录  |  Agenda")

toc = [
    ("01", "大赛总览", "BigAlpha 与三大赛事框架"),
    ("02", "共同信息", "赛程 · 决赛 · 奖项"),
    ("03", "AI 因子挖掘", "传统 + AI 双赛道，挖掘 Alpha 因子"),
    ("04", "端到端大模型", "原始量价数据直出预测"),
    ("05", "AI 开放创新", "围绕 AI × 量化自由立项"),
    ("06", "三赛事对比", "横向差异 + 选题建议"),
]

col_w = Inches(6.2)
row_h = Inches(0.85)
for i, (num, title, desc) in enumerate(toc):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * (col_w + Inches(0.2))
    top = Inches(1.6) + row * (row_h + Inches(0.25))
    add_rect(s, left, top, col_w, row_h, fill=LIGHT)
    add_rect(s, left, top, Inches(0.85), row_h, fill=NAVY)
    add_text(s, left, top, Inches(0.85), row_h, num,
             size=24, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(1.0), top + Inches(0.1),
             col_w - Inches(1.1), Inches(0.4), title,
             size=16, bold=True, color=NAVY)
    add_text(s, left + Inches(1.0), top + Inches(0.45),
             col_w - Inches(1.1), Inches(0.35), desc,
             size=11, color=GREY)

# =============================================================================
# Slide 3 — 大赛总览
# =============================================================================
s = add_slide()
page_header(s, "01  大赛总览  |  BigAlpha 全球联赛",
            "面向全球高校青年人才的年度 AI 量化投研赛事 IP")

# 顶部大字介绍
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
         "连接高校 · 学生 · 量化机构 · AI 金融实践场景",
         size=18, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.7),
         "通过真实数据、开放赛题与专业评审，发现具备 AI 投研能力、量化建模能力与策略创新能力的新一代金融科技人才。",
         size=13, color=GREY)

# 三大赛事卡片
tracks = [
    ("AI 因子挖掘",
     "金融工程 / 量化 / 数学 / 统计",
     "在分钟级 A 股数据上挖掘日频因子，传统赛道与 AI 赛道并行评分。",
     BLUE),
    ("端到端大模型",
     "AI / CS / 数据科学 / 金融科技",
     "无特征工程，用深度网络直接从原始量价序列学习残差收益预测。",
     PURPLE),
    ("AI 开放创新",
     "跨学科 / 创新创业团队",
     "围绕 AI × 量化任意立项：因子、事件驱动、多模态、Agent 工具链。",
     ORANGE),
]

card_w = Inches(4.0)
card_h = Inches(3.5)
gap = Inches(0.15)
total_w = card_w * 3 + gap * 2
start_left = (SW - total_w) // 2
top = Inches(3.5)

for i, (name, audience, desc, color) in enumerate(tracks):
    left = start_left + i * (card_w + gap)
    add_rect(s, left, top, card_w, card_h, fill=LIGHT)
    add_rect(s, left, top, card_w, Inches(0.7), fill=color)
    add_text(s, left, top, card_w, Inches(0.7), name,
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.25), top + Inches(0.95),
             card_w - Inches(0.5), Inches(0.4),
             "适合人群", size=11, bold=True, color=color)
    add_text(s, left + Inches(0.25), top + Inches(1.3),
             card_w - Inches(0.5), Inches(0.6),
             audience, size=12, color=GREY)
    add_text(s, left + Inches(0.25), top + Inches(2.0),
             card_w - Inches(0.5), Inches(0.4),
             "赛题概览", size=11, bold=True, color=color)
    add_text(s, left + Inches(0.25), top + Inches(2.35),
             card_w - Inches(0.5), Inches(1.0),
             desc, size=12, color=GREY)

# =============================================================================
# Slide 4 — BigQuant 平台介绍
# =============================================================================
s = add_slide()
page_header(s, "01  关于 BigQuant  |  AI 量化投研平台",
            "Democratize AI to empower investors —— 让 AI 普惠每一位投资者")

add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.6),
         "国内领先的 AI 量化投资平台，行业内首个将 AI 系统性应用于投资领域的平台级产品；本届赛事直接复用机构客户一致的工具链。",
         size=12, color=GREY)

caps = [
    ("数据底座",
     "PB 级标准金融数据 + 另类数据，覆盖行情、财务、舆情、产业链；提供 PIT 处理、跨频率对齐等基础设施。"),
    ("因子与算法",
     "2000+ 基础因子库；表达式引擎 / UDF / AI 自动化挖掘；AutoML、超参寻优、滚动训练、组合优化、归因分析等组件。"),
    ("研发环境",
     "DAI 数据引擎提供高性能因子计算（速度更快 / 内存更低）；模块化可视化开发与 Python / Notebook 无缝集成。"),
    ("生态与落地",
     "策略源码库、券商研报、量化学院培训、模拟与实盘对接，构建从学习到落地的完整闭环。"),
]
card_w = Inches(6.0)
card_h = Inches(2.1)
top = Inches(2.0)
for i, (h, b) in enumerate(caps):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * (card_w + Inches(0.3))
    t = top + row * (card_h + Inches(0.2))
    add_rect(s, left, t, card_w, card_h, fill=LIGHT)
    add_rect(s, left, t, Inches(0.12), card_h, fill=ACCENT)
    add_text(s, left + Inches(0.3), t + Inches(0.2),
             card_w - Inches(0.5), Inches(0.5),
             h, size=15, bold=True, color=NAVY)
    add_text(s, left + Inches(0.3), t + Inches(0.75),
             card_w - Inches(0.5), card_h - Inches(0.85),
             b, size=12, color=GREY)

# =============================================================================
# Slide 5 — 共同赛程时间轴
# =============================================================================
s = add_slide()
page_header(s, "02  共同信息  |  赛程时间轴", "三大赛事统一节奏：报名 → 初赛 → 决赛")

# 主轴
axis_y = Inches(3.4)
add_rect(s, Inches(0.8), axis_y, Inches(11.7), Inches(0.06), fill=NAVY)

stages = [
    ("阶段一 · 宣传报名",
     "2026-05-15 ~ 06-24",
     "5/25 起正式报名 · 单队最多 5 人\n6/8 ~ 6/18 系统内测，结束后清榜"),
    ("阶段二 · 初赛",
     "2026-06-25 ~ 08-05",
     "因子挖掘 / 端到端 7/26 截止；\n开放创新 8/5 提交方案；\n公榜实时打分，私榜决定排名"),
    ("阶段三 · 决赛",
     "2026-09",
     "美国场 09-01 @ UC Berkeley\n亚洲场 9 月中旬 @ 北京大学\n现场答辩 + 颁奖"),
]

block_w = Inches(3.7)
block_h = Inches(2.2)
for i, (name, when, body) in enumerate(stages):
    cx = Inches(0.8) + Inches(11.7) * (0.166 + 0.333 * i)
    left = cx - block_w // 2
    top = Inches(1.0) if i % 2 == 0 else Inches(4.0)
    add_rect(s, left, top, block_w, block_h, fill=LIGHT)
    add_rect(s, left, top, block_w, Inches(0.55), fill=BLUE)
    add_text(s, left, top, block_w, Inches(0.55), name,
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left, top + Inches(0.6), block_w, Inches(0.4), when,
             size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, left + Inches(0.2), top + Inches(1.0),
             block_w - Inches(0.4), block_h - Inches(1.05),
             body, size=11, color=GREY, align=PP_ALIGN.CENTER)
    # 节点
    dot = add_rect(s, cx - Inches(0.12), axis_y - Inches(0.09),
                   Inches(0.24), Inches(0.24),
                   fill=ACCENT, shape=MSO_SHAPE.OVAL)

# =============================================================================
# Slide 6 — 决赛与晋级规则
# =============================================================================
s = add_slide()
page_header(s, "02  共同信息  |  决赛与晋级规则",
            "两地双场 · 12 个晋级名额 · 现场答辩")

# 左：场次
add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(2.7), fill=LIGHT)
add_text(s, Inches(0.7), Inches(1.35), Inches(5.6), Inches(0.5),
         "决赛场次", size=16, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(0.5),
         "美国场 · 2026-09-01 · UC Berkeley",
         size=14, bold=True, color=BLUE)
add_text(s, Inches(0.7), Inches(2.35), Inches(5.6), Inches(0.5),
         "亚洲场 · 2026 年 9 月中旬 · 北京大学",
         size=14, bold=True, color=BLUE)
add_text(s, Inches(0.7), Inches(3.0), Inches(5.6), Inches(0.8),
         "答辩形式：25 分钟展示 + 5 分钟评委问答；\n线下举行，同步线上直播。",
         size=12, color=GREY)

# 右：晋级规则
add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.7), fill=LIGHT)
add_text(s, Inches(7.0), Inches(1.35), Inches(5.6), Inches(0.5),
         "晋级名额（每场 12 个）", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(7.0), Inches(1.95), Inches(5.6), Inches(2.0), [
    ("直通晋级 · 6 席", "每场 × 各赛事初赛前 2 名（3 × 2）"),
    ("综合评定 · 6 席", "组委会跨赛事综合判定选出"),
    ("评定来源", "三大赛事初赛表现整体评估"),
], size=12)

# 决赛材料
add_rect(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.9), fill=LIGHT)
add_text(s, Inches(0.7), Inches(4.25), Inches(12), Inches(0.4),
         "决赛参赛材料：研究报告（PDF · 10–15 页）",
         size=15, bold=True, color=NAVY)
sections = [
    ("摘要", "研究目标 · 方法 · 结论"),
    ("引言", "研究背景与文献综述"),
    ("因子 / 模型构建", "数据预处理与方法细节"),
    ("实证分析", "回测、归因、稳健性"),
    ("创新与局限", "方法创新点与边界"),
    ("结论", "总结与未来方向"),
]
for i, (h, b) in enumerate(sections):
    col = i % 3
    row = i // 3
    left = Inches(0.7) + col * Inches(4.05)
    top = Inches(4.85) + row * Inches(1.05)
    add_rect(s, left, top, Inches(3.85), Inches(0.95), fill=WHITE,
             line=SOFT_BLUE)
    add_text(s, left + Inches(0.15), top + Inches(0.1),
             Inches(3.55), Inches(0.4), h,
             size=13, bold=True, color=BLUE)
    add_text(s, left + Inches(0.15), top + Inches(0.45),
             Inches(3.55), Inches(0.45), b,
             size=11, color=GREY)

# =============================================================================
# Slide 7 — 共同奖项体系
# =============================================================================
s = add_slide()
page_header(s, "02  共同信息  |  奖项体系",
            "三大赛事共享奖项；覆盖比赛全周期")

# 洲际奖项表头
add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
         "洲际奖项（每场分别设置）", size=16, bold=True, color=NAVY)

prizes = [
    ("金奖", "1 名", "￥20,000", "奖杯 + 电子证书", ACCENT),
    ("银奖", "2 名", "￥15,000", "奖杯 + 电子证书", BLUE),
    ("铜奖", "3 名", "￥10,000", "奖杯 + 电子证书", PURPLE),
]
card_w = Inches(4.0)
card_h = Inches(2.0)
top = Inches(1.85)
gap = Inches(0.15)
total = card_w * 3 + gap * 2
start = (SW - total) // 2
for i, (name, count, money, ext, color) in enumerate(prizes):
    left = start + i * (card_w + gap)
    add_rect(s, left, top, card_w, card_h, fill=LIGHT)
    add_rect(s, left, top, card_w, Inches(0.55), fill=color)
    add_text(s, left, top, card_w, Inches(0.55), name,
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left, top + Inches(0.65), card_w, Inches(0.4), count,
             size=13, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, left, top + Inches(1.0), card_w, Inches(0.5), money,
             size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, left, top + Inches(1.55), card_w, Inches(0.4), ext,
             size=11, color=GREY, align=PP_ALIGN.CENTER)

# 初赛 + 特色奖
add_rect(s, Inches(0.5), Inches(4.2), Inches(6.0), Inches(2.7), fill=LIGHT)
add_text(s, Inches(0.7), Inches(4.35), Inches(5.6), Inches(0.4),
         "初赛奖项", size=15, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(4.85), Inches(5.6), Inches(2.0), [
    ("赛道周冠军", "每周评选 · ￥1,500 / 周"),
    ("BigAlpha 量化新星", "至少提交一次有效作品即颁发电子证书"),
], size=12)

add_rect(s, Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.7), fill=LIGHT)
add_text(s, Inches(7.0), Inches(4.35), Inches(5.6), Inches(0.4),
         "特色奖项", size=15, bold=True, color=NAVY)
add_text(s, Inches(7.0), Inches(4.85), Inches(5.6), Inches(2),
         "随赛事推进，组委会评选具有特色的作品与团队进行专项颁奖；奖项内容届时公布，奖金合计 ￥12,500。",
         size=12, color=GREY)


# -----------------------------------------------------------------------------
# 通用：赛道章节封面
# -----------------------------------------------------------------------------
def section_cover(num, name_zh, name_en, theme, tagline):
    s = add_slide()
    add_rect(s, 0, 0, SW, SH, fill=NAVY)
    add_rect(s, 0, Inches(4.6), SW, Inches(0.05), fill=ACCENT)
    add_text(s, Inches(0.8), Inches(1.6), Inches(12), Inches(0.5),
             num, size=22, bold=True, color=ACCENT)
    add_text(s, Inches(0.8), Inches(2.1), Inches(12), Inches(1.2),
             name_zh, size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.2), Inches(12), Inches(0.6),
             name_en, size=20, color=SOFT_BLUE)
    add_text(s, Inches(0.8), Inches(4.9), Inches(12), Inches(0.6),
             theme, size=18, color=ACCENT, bold=True)
    add_text(s, Inches(0.8), Inches(5.5), Inches(12), Inches(1.2),
             tagline, size=14, color=SOFT_BLUE)
    return s


# =============================================================================
# Slide 8 — 因子挖掘 章节封面
# =============================================================================
section_cover(
    "03",
    "AI 因子挖掘",
    "Alpha Factor Mining",
    "传统赛道 + AI 智能赛道并行",
    "在中证 1000 成分股的分钟级行情上构建日频因子；公榜实时迭代、私榜样本外评分。\n适合：金融工程、量化、数学、统计方向学生。",
)

# =============================================================================
# Slide 9 — 因子挖掘：赛道与数据
# =============================================================================
s = add_slide()
page_header(s, "03  AI 因子挖掘  |  双赛道与数据", "投资逻辑 vs AI 创新；同一数据集，不同评审导向")

# 双赛道
add_text(s, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
         "双赛道并行模式", size=16, bold=True, color=NAVY)

tracks = [
    ("传统量化赛道", BLUE,
     "强调投资逻辑、统计显著性与经济解释性。",
     ["多因子模型 · 基本面 · 经典技术指标 · 物理模型",
      "重点评估：经济逻辑合理性、统计显著性、跨市场环境稳健性、可解释性。"]),
    ("AI 智能赛道", PURPLE,
     "利用 LLM、强化学习、遗传算法做因子自动生成或优化。",
     ["自动化特征工程 · Prompt 生成因子公式 · 神经网络挖掘非线性",
      "重点评估：因子有效性 + AI 应用深度 / 创新性 / 参与度。"]),
]
top = Inches(1.6)
card_w = Inches(6.1)
card_h = Inches(2.7)
for i, (name, color, sub, body) in enumerate(tracks):
    left = Inches(0.5) + i * (card_w + Inches(0.2))
    add_rect(s, left, top, card_w, card_h, fill=LIGHT)
    add_rect(s, left, top, card_w, Inches(0.5), fill=color)
    add_text(s, left, top, card_w, Inches(0.5), name,
             size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.25), top + Inches(0.6),
             card_w - Inches(0.5), Inches(0.4), sub,
             size=12, bold=True, color=color)
    add_bullets(s, left + Inches(0.25), top + Inches(1.05),
                card_w - Inches(0.5), card_h - Inches(1.1),
                body, size=11)

# 数据
add_rect(s, Inches(0.5), Inches(4.55), Inches(12.3), Inches(2.55), fill=LIGHT)
add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4),
         "数据  |  高质量 A 股市场数据", size=15, bold=True, color=NAVY)

data_items = [
    ("股票池", "中证 1000 指数历史相应时点成分股"),
    ("时间范围", "2019-01-01 ~ 2024-12-31"),
    ("行情数据", "1 分钟 K 线 + 盘口快照"),
    ("财务数据", "已做 PIT（Point-in-Time）处理"),
]
for i, (h, b) in enumerate(data_items):
    col = i % 2
    row = i // 2
    left = Inches(0.7) + col * Inches(6.1)
    top = Inches(5.2) + row * Inches(0.85)
    add_rect(s, left, top, Inches(0.15), Inches(0.7), fill=ACCENT)
    add_text(s, left + Inches(0.3), top, Inches(2.0), Inches(0.35),
             h, size=12, bold=True, color=NAVY)
    add_text(s, left + Inches(0.3), top + Inches(0.35),
             Inches(5.6), Inches(0.4), b, size=12, color=GREY)

# =============================================================================
# Slide 10 — 因子挖掘：评估体系
# =============================================================================
s = add_slide()
page_header(s, "03  AI 因子挖掘  |  评估体系",
            "团队最终得分 = 0.3 × A 项（最佳单因子） + 0.7 × B 项（Elastic Net 组合贡献）")

# A 项
add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(3.5), fill=LIGHT)
add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.5), fill=BLUE)
add_text(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.5),
         "A 项 · 单因子最佳得分（30%）", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(0.5),
         "FACTOR = 0.25·Rank(IC_mean) + 0.25·Rank(IC_IR)",
         size=12, color=NAVY, bold=True)
add_text(s, Inches(0.7), Inches(2.2), Inches(5.6), Inches(0.5),
         "          + 0.25·Rank(SR) + 0.25·Rank(Stress)",
         size=12, color=NAVY, bold=True)
add_bullets(s, Inches(0.7), Inches(2.7), Inches(5.6), Inches(2.0), [
    ("IC 均值", "截面相关性的强度"),
    ("IC_IR", "IC 序列稳定性"),
    ("SR", "多空组合年化夏普"),
    ("Stress", "特殊行情下的稳健性"),
], size=11)

# B 项
add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(3.5), fill=LIGHT)
add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.5), fill=PURPLE)
add_text(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.5),
         "B 项 · Elastic Net 组合贡献（70%）",
         size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(7.0), Inches(1.85), Inches(5.6), Inches(2.8), [
    ("汇总全场", "所有团队入池因子构成全局候选集"),
    ("滚动回归", "60 个交易日窗口、步长 20 日"),
    ("ModelScore", "mean(|w|) / (std(|w|) + ε)"),
    ("L1 选因子", "无增量贡献者权重压至 0"),
    ("团队得分", "本队被选中因子的 ModelScore 百分位均值"),
], size=11)

# 提交规则 + 数据划分
add_rect(s, Inches(0.5), Inches(4.95), Inches(6.0), Inches(2.05), fill=LIGHT)
add_text(s, Inches(0.7), Inches(5.05), Inches(5.6), Inches(0.4),
         "提交规则", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(5.4), Inches(5.6), Inches(1.6), [
    "每队最多提交 50 个因子",
    "从中选 1 ~ 50 个进入因子池参评",
    "鼓励精选，多堆砌反而拉低 B 项",
], size=11)

add_rect(s, Inches(6.8), Inches(4.95), Inches(6.0), Inches(2.05), fill=LIGHT)
add_text(s, Inches(7.0), Inches(5.05), Inches(5.6), Inches(0.4),
         "数据划分（公榜 / 私榜）", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(7.0), Inches(5.4), Inches(5.6), Inches(1.6), [
    ("公榜", "训练 2019-01-01 ~ 2024-12-31 / 验证 2025 全年"),
    ("私榜", "区间不公开，含 2026 样本外，最终排名以此为准"),
], size=11)

# =============================================================================
# Slide 11 — 因子挖掘：提交与代码要求
# =============================================================================
s = add_slide()
page_header(s, "03  AI 因子挖掘  |  提交规范",
            "提交因子构建代码（非数据），平台自动运行评分")

# 数据格式
add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(3.0), fill=LIGHT)
add_text(s, Inches(0.7), Inches(1.35), Inches(5.6), Inches(0.4),
         "返回数据格式（main 函数输出）", size=14, bold=True, color=NAVY)

cols = ["date", "instrument", "factor"]
rows = [
    ["2023-01-03 09:45:00", "000001.SZ", "0.05"],
    ["2023-01-03 09:45:00", "000002.SZ", "-0.12"],
    ["...", "...", "..."],
]
tab_left = Inches(0.7)
tab_top = Inches(1.85)
col_w = Inches(1.85)
row_h = Inches(0.45)
# header
for j, c in enumerate(cols):
    add_rect(s, tab_left + j * col_w, tab_top, col_w, row_h, fill=NAVY)
    add_text(s, tab_left + j * col_w, tab_top, col_w, row_h, c,
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        add_rect(s, tab_left + j * col_w, tab_top + (i + 1) * row_h,
                 col_w, row_h, fill=WHITE, line=SOFT_BLUE)
        add_text(s, tab_left + j * col_w, tab_top + (i + 1) * row_h,
                 col_w, row_h, val, size=10, color=GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.7), Inches(3.5), Inches(5.6), Inches(0.7),
         "默认因子值越大越好；参赛者需自行确保因子方向正确性。",
         size=11, color=GREY)

# 数据校验
add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(3.0), fill=LIGHT)
add_text(s, Inches(7.0), Inches(1.35), Inches(5.6), Inches(0.4),
         "数据校验与预处理", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(7.0), Inches(1.8), Inches(5.6), Inches(2.6), [
    ("列检查", "仅含 date / instrument / factor 三列"),
    ("交易日完整", "区间内不缺任何交易日"),
    ("覆盖度", "每日缺失率 ≤ 40%"),
    ("预处理", "去极值 + 标准化"),
    ("风格剔除", "对 BARRA 风险因子回归取残差"),
], size=11)

# 代码要求
add_rect(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.6), fill=LIGHT)
add_text(s, Inches(0.7), Inches(4.55), Inches(12), Inches(0.4),
         "平台与代码要求", size=14, bold=True, color=NAVY)
items = [
    ("平台运行", "在指定线上 Notebook 环境运行；AI 赛道使用平台 AI 工具链 / 算力"),
    ("运行时长", "CPU Notebook ≤ 3 小时；AI 赛道 GPU ≤ 6 小时"),
    ("网络限制", "Notebook 互联网访问被禁用，防止信息泄露与未来数据"),
    ("最终提交", "提交可自动运行生成因子的代码；AI 赛道额外提交应用说明"),
]
for i, (h, b) in enumerate(items):
    col = i % 2
    row = i // 2
    left = Inches(0.7) + col * Inches(6.1)
    top = Inches(5.05) + row * Inches(0.95)
    add_text(s, left, top, Inches(1.5), Inches(0.4), h,
             size=12, bold=True, color=BLUE)
    add_text(s, left + Inches(1.6), top, Inches(4.4), Inches(0.85),
             b, size=11, color=GREY)

# =============================================================================
# Slide 12 — 端到端模型 章节封面
# =============================================================================
section_cover(
    "04",
    "端到端大模型",
    "End-to-End Quant Modeling",
    "原始量价 → 投资决策一站式映射",
    "无显式特征工程，用深度网络直接从分钟级行情学习残差收益预测；私榜阶段平台重训。\n适合：AI、CS、数据科学、金融科技方向学生。",
)

# =============================================================================
# Slide 13 — 端到端：任务与数据
# =============================================================================
s = add_slide()
page_header(s, "04  端到端大模型  |  任务与数据",
            "鼓励无特征 / 弱特征建模，直接处理原始量价序列")

# 任务概览
add_rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.4), fill=LIGHT)
add_rect(s, Inches(0.5), Inches(1.2), Inches(0.15), Inches(2.4), fill=ACCENT)
add_text(s, Inches(0.85), Inches(1.35), Inches(12), Inches(0.4),
         "任务定义", size=15, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(1.85), Inches(12), Inches(1.7),
         "传统 “因子挖掘 + 组合优化” 两步走范式存在信息损耗与人工偏见；\n"
         "本赛道聚焦端到端模型生成：用多频率分钟级行情作为输入，训练目标为股票残差收益率，\n"
         "通过模型内部机制自发捕捉市场微观结构的深层规律。",
         size=12, color=GREY)

# 数据卡
add_text(s, Inches(0.5), Inches(3.85), Inches(12), Inches(0.4),
         "数据说明", size=15, bold=True, color=NAVY)
data_cards = [
    ("股票池", "中证 1000 历史成分股", BLUE),
    ("时间范围", "2019-01-01 ~ 2023-12-31", PURPLE),
    ("行情频率", "1 / 5 / 15 / 30 分钟 K 线", ORANGE),
    ("市场快照", "盘口数据（多频率）", GREEN),
]
card_w = Inches(3.0)
card_h = Inches(1.7)
gap = Inches(0.1)
total = card_w * 4 + gap * 3
start = (SW - total) // 2
top = Inches(4.35)
for i, (h, b, c) in enumerate(data_cards):
    left = start + i * (card_w + gap)
    add_rect(s, left, top, card_w, card_h, fill=LIGHT)
    add_rect(s, left, top, card_w, Inches(0.45), fill=c)
    add_text(s, left, top, card_w, Inches(0.45), h,
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left, top + Inches(0.55), card_w, Inches(1.05), b,
             size=12, color=GREY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 模板代码
add_text(s, Inches(0.5), Inches(6.25), Inches(12), Inches(0.4),
         "官方模板代码", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.7),
         "Transformer 基础模型（demo_transformer.ipynb）  ·  轻量化序列模型（demo_lightweight.ipynb）  ·  分布式数据加载（demo_data_loader.ipynb）",
         size=11, color=GREY)

# =============================================================================
# Slide 14 — 端到端：模型规范
# =============================================================================
s = add_slide()
page_header(s, "04  端到端大模型  |  模型规范",
            "硬性约束限制特征工程，鼓励真正的端到端结构")

# 输入约束
add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(3.5), fill=LIGHT)
add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.5), fill=BLUE)
add_text(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.5),
         "输入与特征约束", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(2.85), [
    ("特征上限", "≤ 100 个原始字段（时序长度不计）"),
    ("数据来源", "全部来自主办方提供的数据"),
    ("禁止操作", "跨字段算子、滚动统计、因子合成、降维、第三方数据"),
    ("允许预处理", "缺失填充、归一化、对数 / 符号变换（仅训练集统计）"),
    ("时序窗口", "单次推理回看 ≤ 240 个交易日"),
], size=11)

# 模型架构
add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(3.5), fill=LIGHT)
add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.5), fill=PURPLE)
add_text(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.5),
         "模型架构约束", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(7.0), Inches(1.85), Inches(5.6), Inches(2.85), [
    ("可训练参数", "10 万 ≤ 参数量 ≤ 1 亿"),
    ("外部权重", "禁止使用任何外部预训练权重"),
    ("从零训练", "所有参数基于本竞赛数据训练"),
    ("架构形态", "未硬性指定；浅层模型（线性 / XGBoost）难有竞争力"),
], size=11)

# 输出格式
add_rect(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(2.0), fill=LIGHT)
add_text(s, Inches(0.7), Inches(5.05), Inches(5.6), Inches(0.4),
         "模型输出格式（推理产出分数文件）",
         size=14, bold=True, color=NAVY)
cols = ["date", "instrument", "score"]
rows = [["2023-01-03 09:45:00", "000001.SZ", "0.05"],
        ["2023-01-03 09:45:00", "000002.SZ", "-0.12"]]
tab_left = Inches(0.7)
tab_top = Inches(5.5)
col_w = Inches(2.6)
row_h = Inches(0.4)
for j, c in enumerate(cols):
    add_rect(s, tab_left + j * col_w, tab_top, col_w, row_h, fill=NAVY)
    add_text(s, tab_left + j * col_w, tab_top, col_w, row_h, c,
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        add_rect(s, tab_left + j * col_w, tab_top + (i + 1) * row_h,
                 col_w, row_h, fill=WHITE, line=SOFT_BLUE)
        add_text(s, tab_left + j * col_w, tab_top + (i + 1) * row_h,
                 col_w, row_h, val, size=10, color=GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(8.6), Inches(5.5), Inches(4.0), Inches(1.5),
         "默认分数越大越好；\n所产出分数经风格剔除后等价于每日更新的单因子，\n因此直接以单因子得分作为团队得分。",
         size=11, color=GREY)

# =============================================================================
# Slide 15 — 端到端：评估与提交
# =============================================================================
s = add_slide()
page_header(s, "04  端到端大模型  |  评估与提交",
            "公榜按提交权重推理；私榜由平台从零重训后推理")

# 最终得分公式
add_rect(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(1.05), fill=LIGHT)
add_text(s, Inches(0.7), Inches(1.25), Inches(12), Inches(0.4),
         "最终得分公式（与因子挖掘 A 项一致）",
         size=13, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.5),
         "Score = 0.25·Rank(IC_mean) + 0.25·Rank(IC_IR) + 0.25·Rank(SR) + 0.25·Rank(Stress)",
         size=14, bold=True, color=BLUE)

# 平台预处理
add_rect(s, Inches(0.5), Inches(2.35), Inches(6.0), Inches(2.4), fill=LIGHT)
add_text(s, Inches(0.7), Inches(2.45), Inches(5.6), Inches(0.4),
         "平台预处理流程", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(2.9), Inches(5.6), Inches(1.8), [
    ("去极值", "截面 1% / 99% 分位 winsorize"),
    ("标准化", "截面 z-score"),
    ("风格剔除", "对 10 类 BARRA 风险因子回归取残差"),
], size=11)

# 数据划分
add_rect(s, Inches(6.8), Inches(2.35), Inches(6.0), Inches(2.4), fill=LIGHT)
add_text(s, Inches(7.0), Inches(2.45), Inches(5.6), Inches(0.4),
         "数据划分", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(7.0), Inches(2.9), Inches(5.6), Inches(1.8), [
    ("公榜训练", "2019-01-01 ~ 2023-12-31"),
    ("公榜验证", "2024-01-01 ~ 2024-12-31"),
    ("私榜", "区间不公开，含样本外数据，最终排名"),
], size=11)

# 提交规范
add_rect(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(2.05), fill=LIGHT)
add_text(s, Inches(0.7), Inches(5.05), Inches(12), Inches(0.4),
         "提交规范与频次", size=14, bold=True, color=NAVY)
items = [
    ("公榜单次", "完整模型包：训练 + 推理脚本、依赖、超参、随机种子、已训练权重"),
    ("公榜推理", "平台仅基于提交权重做推理，免去重复训练成本"),
    ("提交频次", "每队每日 ≤ 3 次（含失败提交）；提交互不覆盖"),
    ("私榜模型", "公榜截止前指定唯一参赛模型，逾期默认取最高分提交并冻结"),
    ("私榜重训", "平台隔离环境从零重训 + 私榜区间推理，结果即为最终排名"),
]
for i, (h, b) in enumerate(items):
    row = i // 1
    left = Inches(0.7)
    top = Inches(5.45) + i * Inches(0.32)
    add_text(s, left, top, Inches(1.7), Inches(0.3), h,
             size=11, bold=True, color=BLUE)
    add_text(s, left + Inches(1.7), top, Inches(10.3), Inches(0.3),
             b, size=11, color=GREY)

# =============================================================================
# Slide 16 — AI 开放创新 章节封面
# =============================================================================
section_cover(
    "05",
    "AI 开放创新",
    "Open Innovation in AI × Quant",
    "面向 AI × 量化的开放式创新创业项目",
    "不预设固定命题、不逐日打分；自由立项，线下路演评比，配套人才网络与孵化资源。\n适合：跨学科、创新创业团队。",
)

# =============================================================================
# Slide 17 — 创新赛：方向与数据
# =============================================================================
s = add_slide()
page_header(s, "05  AI 开放创新  |  立项方向与数据工具栈",
            "鼓励原创 · 可跨方向融合 · AI 必须为核心")

add_text(s, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
         "可选立项方向（提交时注明主方向即可）",
         size=15, bold=True, color=NAVY)

dirs = [
    ("AI 因子与策略", "LLM、时序模型、强化学习挖掘特色因子；端到端 AI 量化策略", BLUE),
    ("宏观与事件驱动", "美联储议息、地缘政治、政策调整：风险识别 / 情景模拟 / 对冲", PURPLE),
    ("多模态金融理解", "研报、电话会、新闻、图表的多模态解析与信号生成", ORANGE),
    ("AI 工具与基础设施", "评测框架、回测加速、可解释性、Agent 工作流等", GREEN),
]
top = Inches(1.6)
card_w = Inches(6.1)
card_h = Inches(1.4)
for i, (h, b, c) in enumerate(dirs):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * (card_w + Inches(0.2))
    t = top + row * (card_h + Inches(0.2))
    add_rect(s, left, t, card_w, card_h, fill=LIGHT)
    add_rect(s, left, t, Inches(0.15), card_h, fill=c)
    add_text(s, left + Inches(0.3), t + Inches(0.15),
             card_w - Inches(0.5), Inches(0.45), h,
             size=14, bold=True, color=NAVY)
    add_text(s, left + Inches(0.3), t + Inches(0.65),
             card_w - Inches(0.5), Inches(0.7), b,
             size=11, color=GREY)

# 数据与工具
add_rect(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.2), fill=LIGHT)
add_text(s, Inches(0.7), Inches(4.95), Inches(12), Inches(0.4),
         "数据与工具栈", size=14, bold=True, color=NAVY)
data_items = [
    ("结构化数据", "A 股日频 / 分钟频行情、PIT 财务数据、宏观指标"),
    ("非结构化数据", "研报、新闻、电话会、股吧舆情（已合规脱敏）"),
    ("算力", "CPU / GPU Notebook；决赛阶段为入围团队提供 GPU 算力包"),
    ("AI 工具链", "主流开源 LLM 接口 + BigQuant 自研量化基础模型 API"),
]
for i, (h, b) in enumerate(data_items):
    col = i % 2
    row = i // 2
    left = Inches(0.7) + col * Inches(6.1)
    top = Inches(5.45) + row * Inches(0.7)
    add_text(s, left, top, Inches(1.7), Inches(0.3),
             h, size=11, bold=True, color=BLUE)
    add_text(s, left + Inches(1.7), top, Inches(4.3), Inches(0.6),
             b, size=11, color=GREY)

# =============================================================================
# Slide 18 — 创新赛：准入与评审
# =============================================================================
s = add_slide()
page_header(s, "05  AI 开放创新  |  准入校验与决赛评审",
            "通过准入 → 进入决赛评审池 → 路演评分")

# 准入校验
add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(3.6), fill=LIGHT)
add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.5), fill=ORANGE)
add_text(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.5),
         "方案准入校验", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(2.95), [
    ("内容相关性", "与量化领域强相关；不违反金融监管或公序良俗"),
    ("AI 核心性", "核心逻辑必须由 AI 主导，禁止贴标签式 AI"),
    ("材料完整性", "Notebook + 方案说明 + AI 应用说明 三件套齐全"),
    ("合规与安全", "无外部网络访问；LLM 输出附事实核查与免责声明"),
], size=11)

# 决赛评审
add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(3.6), fill=LIGHT)
add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.5), fill=PURPLE)
add_text(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.5),
         "决赛评审标准（百分制）", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
items = [
    ("AI 应用创新性 · 35%", "技术选型适配性、原创性"),
    ("方案完整性 · 25%", "数据 / 训练 / 验证流程严谨；可复现"),
    ("商业化与生态 · 20%", "客户画像、付费意愿、与 BigQuant 协同"),
    ("现场表现 · 20%", "陈述清晰、问答精准"),
]
for i, (h, b) in enumerate(items):
    top = Inches(1.95) + i * Inches(0.7)
    add_text(s, Inches(7.0), top, Inches(2.6), Inches(0.3),
             h, size=11, bold=True, color=NAVY)
    add_text(s, Inches(7.0), top + Inches(0.3),
             Inches(5.6), Inches(0.4), b, size=11, color=GREY)

# 提交三件套
add_rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), fill=LIGHT)
add_text(s, Inches(0.7), Inches(5.1), Inches(12), Inches(0.4),
         "提交三件套（方案研发阶段）",
         size=14, bold=True, color=NAVY)
items = [
    ("核心 Notebook",
     "main 函数承载核心逻辑；AI 关键环节添加 # [AI-CORE] 注释"),
    ("方案说明文档",
     "嵌入 Notebook：场景定义 / 数据使用 / AI 方案 / 验证 / 商业化假设"),
    ("AI 应用说明表",
     "逐项列出 AI 介入位置、所用模型、输入输出、可替代性评估"),
]
for i, (h, b) in enumerate(items):
    col = i
    left = Inches(0.7) + col * Inches(4.05)
    top = Inches(5.55)
    add_rect(s, left, top, Inches(3.85), Inches(1.3), fill=WHITE,
             line=SOFT_BLUE)
    add_text(s, left + Inches(0.15), top + Inches(0.1),
             Inches(3.55), Inches(0.4), h, size=12, bold=True, color=BLUE)
    add_text(s, left + Inches(0.15), top + Inches(0.5),
             Inches(3.55), Inches(0.8), b, size=11, color=GREY)

# =============================================================================
# Slide 19 — 三大赛事横向对比
# =============================================================================
s = add_slide()
page_header(s, "06  三大赛事对比  |  快速选题指南",
            "同一平台、同一时间窗，但任务形态、评分逻辑差异显著")

cols = ["维度", "AI 因子挖掘", "端到端大模型", "AI 开放创新"]
rows = [
    ["产出物", "日频 alpha 因子", "端到端模型 + 权重", "完整 AI × 量化方案"],
    ["数据", "1 分钟 K 线 + PIT 财务\n2019 ~ 2024", "1/5/15/30 分钟行情\n2019 ~ 2023", "结构化 + 非结构化\n含舆情语料"],
    ["核心约束", "因子方向、覆盖度、风格剔除", "≤100 字段、参数 1e5~1e8\n禁特征工程", "AI 必须为核心\n禁贴标签式"],
    ["评分逻辑", "0.3·单因子最佳 + 0.7·组合贡献", "IC / IC_IR / SR / Stress 等权", "评委百分制（创新 35%）"],
    ["榜单机制", "公榜 + 私榜（自动迭代）", "公榜 + 私榜（平台重训）", "无榜单 · 路演评比"],
    ["截止日", "公榜 7/26 · 决赛 9 月", "公榜 7/26 · 决赛 9 月", "方案 8/5 · 决赛 9 月"],
    ["适合人群", "金工 / 量化 / 数学 / 统计", "AI / CS / 数据科学", "跨学科 / 创新创业"],
]

tab_left = Inches(0.5)
tab_top = Inches(1.2)
total_w = Inches(12.3)
col_widths = [Inches(1.6), Inches(3.6), Inches(3.55), Inches(3.55)]
row_h = Inches(0.7)
header_h = Inches(0.5)
# header
x = tab_left
for j, c in enumerate(cols):
    add_rect(s, x, tab_top, col_widths[j], header_h,
             fill=NAVY)
    add_text(s, x, tab_top, col_widths[j], header_h, c,
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += col_widths[j]
# rows
for i, row in enumerate(rows):
    x = tab_left
    fill = WHITE if i % 2 == 0 else LIGHT
    for j, val in enumerate(row):
        add_rect(s, x, tab_top + header_h + i * row_h,
                 col_widths[j], row_h, fill=fill, line=SOFT_BLUE)
        bold = (j == 0)
        color = NAVY if j == 0 else GREY
        add_text(s, x + Inches(0.1),
                 tab_top + header_h + i * row_h,
                 col_widths[j] - Inches(0.2), row_h,
                 val, size=10, bold=bold, color=color,
                 anchor=MSO_ANCHOR.MIDDLE)
        x += col_widths[j]

# =============================================================================
# Slide 20 — 如何参赛
# =============================================================================
s = add_slide()
page_header(s, "如何参赛  |  报名指引")

add_text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         "参赛流程", size=18, bold=True, color=NAVY)

steps = [
    ("01", "选定赛事", "三大赛事任选其一，可同时参与（账号统一）"),
    ("02", "组队报名", "活动主页报名 · 单队最多 5 人 · 报名截止 6/24"),
    ("03", "加入社群", "报名后加入官方微信 / QQ 群，寻找队友 / 获取通知"),
    ("04", "系统内测", "6/8 ~ 6/18 可提交测试，结束后清榜"),
    ("05", "提交作品", "因子 / 端到端：7/26 截止；开放创新：8/5 截止"),
    ("06", "晋级决赛", "组委会综合评定 12 个晋级名额 · 9 月线下答辩"),
]

start_top = Inches(1.7)
step_w = Inches(12.3)
step_h = Inches(0.7)
for i, (n, h, b) in enumerate(steps):
    top = start_top + i * (step_h + Inches(0.1))
    add_rect(s, Inches(0.5), top, step_w, step_h, fill=LIGHT)
    add_rect(s, Inches(0.5), top, Inches(0.85), step_h, fill=NAVY)
    add_text(s, Inches(0.5), top, Inches(0.85), step_h, n,
             size=18, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.55), top, Inches(2.3), step_h,
             h, size=14, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.0), top, Inches(8.5), step_h,
             b, size=12, color=GREY,
             anchor=MSO_ANCHOR.MIDDLE)

# 报名链接
add_rect(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.45), fill=NAVY)
add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.45),
         "报名入口：https://bigquant.com/square/competition/",
         size=12, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =============================================================================
# 保存
# =============================================================================
out = "/Users/xiehao/Desktop/workspace/BigQuant/BigAlpha/docs/BigAlpha2026_三大赛事介绍.pptx"
prs.save(out)
print(f"saved: {out}  slides={len(prs.slides)}")




